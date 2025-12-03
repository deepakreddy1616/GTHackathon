import os
import openai
from openai import OpenAI
client = OpenAI()

openai.api_key = os.getenv("OPENAI_API_KEY")

# insightflow_run.py
# InsightFlow - Simple Automated Report Maker

import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path
import sys
import subprocess


# Paths and settings
CSV_PATH = "sample_data.csv"
OUTPUT_DIR = Path("output")
CHART_DIR = OUTPUT_DIR / "charts"
PPTX_PATH = OUTPUT_DIR / "InsightFlow_Report.pptx"

DATE_COL = "timestamp"
VALUE_COL = "clicks"

# ---- Step 1: Load the CSV file ----
def load_data(path):
    print("Step 1: Loading data...")
    df = pd.read_csv(path)
    df[DATE_COL] = pd.to_datetime(df[DATE_COL], errors='coerce')
    return df

# ---- Step 2: Clean column names ----
def clean_data(df):
    print("Step 2: Cleaning data...")
    df.columns = [c.strip().lower() for c in df.columns]
    return df

# ---- Step 3: Compute KPIs ----
def compute_kpis(df):
    print("Step 3: Computing KPIs...")
    clicks = int(df["clicks"].sum())
    impressions = int(df["impressions"].sum())
    revenue = float(df["revenue"].sum())
    ctr = clicks / impressions if impressions > 0 else 0
    return clicks, impressions, revenue, ctr

# ---- Step 4: Create chart ----
def create_chart(df):
    print("Step 4: Creating chart...")
    df_sorted = df.sort_values(by=DATE_COL)

    plt.figure(figsize=(8,4))
    plt.plot(df_sorted[DATE_COL], df_sorted["clicks"], marker='o')
    # Draw red dots for weird days
    weird = df_sorted[df_sorted["anomaly"] != "normal"]

    plt.scatter(
    weird[DATE_COL],
    weird["clicks"],
    color="red",
    s=120,
    zorder=5,
)
    plt.title("Daily Clicks")
    plt.xlabel("Date")
    plt.ylabel("Clicks")
    plt.grid(True)

    CHART_DIR.mkdir(parents=True, exist_ok=True)
    chart_path = CHART_DIR / "daily_clicks.png"
    plt.savefig(chart_path)
    plt.close()
    return chart_path
def create_ctr_chart(df):
    print("Creating CTR chart...")
    df_sorted = df.sort_values(by=DATE_COL)

    plt.figure(figsize=(8,4))
    plt.plot(df_sorted[DATE_COL], df_sorted["ctr"], marker='o')
    plt.title("CTR Trend")
    plt.xlabel("Date")
    plt.ylabel("CTR")
    plt.grid(True)

    CHART_DIR.mkdir(parents=True, exist_ok=True)
    ctr_path = CHART_DIR / "ctr_trend.png"
    plt.tight_layout()
    plt.savefig(ctr_path)
    plt.close()
    return ctr_path

def create_revenue_chart(df):
    print("Creating Revenue chart...")
    df_sorted = df.sort_values(by=DATE_COL)

    plt.figure(figsize=(8,4))
    plt.plot(df_sorted[DATE_COL], df_sorted["revenue"], marker='o')
    plt.title("Revenue Trend")
    plt.xlabel("Date")
    plt.ylabel("Revenue")
    plt.grid(True)

    CHART_DIR.mkdir(parents=True, exist_ok=True)
    rev_path = CHART_DIR / "revenue_trend.png"
    plt.tight_layout()
    plt.savefig(rev_path)
    plt.close()
    return rev_path


# ---- Step 5: Create summary text ----
def create_summary(clicks, impressions, revenue, ctr, df):
    print("Step 5: Creating summary text...")
    latest = df.sort_values(by=DATE_COL).iloc[-1]
    latest_date = str(latest[DATE_COL]).split(" ")[0]
    latest_clicks = int(latest["clicks"])

    summary = (
        f"Total Clicks: {clicks}\n"
        f"Total Impressions: {impressions}\n"
        f"Total Revenue: ${revenue:.2f}\n"
        f"CTR: {ctr:.2%}\n"
        f"Latest Day ({latest_date}) Clicks: {latest_clicks}\n\n"
        f"Recommendation: Focus on the days with highest clicks to maximize performance."
    )
    return summary
# ---- AI Summary Enhancement ----
# ---- AI Summary Enhancement (NEW OPENAI CLIENT) ----
def ai_refine_summary(basic_summary):
    """
    Creates a clean AI-enhanced summary using the new OpenAI client (2024+).
    If AI is disabled or fails, returns the original summary.
    """

    # If AI disabled or key missing
    if os.getenv("DISABLE_AI", "") == "1" or not os.getenv("OPENAI_API_KEY"):
        print("AI disabled or no API key found. Using basic summary.")
        return basic_summary

    prompt = (
        "Rewrite this summary into 2–3 clean, simple, professional sentences "
        "and add one meaningful recommendation:\n\n"
        + basic_summary
    )

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "user", "content": prompt}
            ],
            max_tokens=150,
            temperature=0.3,
        )

        ai_text = response.choices[0].message.content.strip()
        return ai_text

    except Exception as e:
        print("\nAI Error:", e)
        print("Using basic summary instead.\n")
        return basic_summary

    # ---- Anomaly detection (simple z-score) ----
def detect_anomalies(ts_df, z_threshold=2.0):
    """
    Input: ts_df with columns ['timestamp','clicks'] (timestamp is datetime)
    Output: same dataframe with a new column 'anomaly' = 'high'/'low'/'normal'
    """
    import numpy as np
    ts = ts_df.copy()
    vals = ts['clicks'].astype(float)
    mean = vals.mean()
    std = vals.std(ddof=0) if vals.std(ddof=0) != 0 else 1e-9
    z = (vals - mean) / std
    ts['anomaly'] = ['high' if zi > z_threshold else ('low' if zi < -z_threshold else 'normal') for zi in z]
    return ts



# ---- Step 6: Create the PowerPoint ----
# ---- Step 6: Create the PowerPoint ----
def create_pptx(summary, chart_path, ctr_chart=None, revenue_chart=None):
    print("Step 6: Creating PowerPoint report...")
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "InsightFlow - Automated Insights Report"

    # Summary slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Executive Summary"
    slide.shapes.placeholders[1].text = summary

    # Chart slide (Daily Clicks)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Daily Clicks Chart"
    slide.shapes.add_picture(str(chart_path), Inches(1), Inches(1.5), width=Inches(8))

    # CTR Chart Slide (if available)
    if ctr_chart is not None:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "CTR Trend"
        slide.shapes.add_picture(str(ctr_chart), Inches(1), Inches(1.5), width=Inches(8))

    # Revenue Chart Slide (if available)
    if revenue_chart is not None:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Revenue Trend"
        slide.shapes.add_picture(str(revenue_chart), Inches(1), Inches(1.5), width=Inches(8))

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)

    print("Report saved at:", PPTX_PATH)
    pptx_to_pdf(PPTX_PATH)

    import subprocess

def pptx_to_pdf(pptx_path):
    p = Path(pptx_path)
    if not p.exists():
        print("pptx_to_pdf: PPTX not found:", p)
        return

    # Default command name -- if not on PATH, replace with full path to soffice.exe
    LO_CMD = r"C:\Program Files\LibreOffice\program\soffice.exe"
  # or r"C:\Program Files\LibreOffice\program\soffice.exe"

    try:
        subprocess.run([LO_CMD, "--headless", "--convert-to", "pdf", str(p), "--outdir", str(p.parent)], check=True)
        print("PDF created at:", p.with_suffix(".pdf"))
    except Exception as e:
        print("PDF conversion error:", e)
        print("If LibreOffice is not in PATH, set LO_CMD to full soffice executable path.")



# ---- Main flow ----
def main():
    if not Path(CSV_PATH).exists():
        print("sample_data.csv not found!")
        sys.exit(1)

    df = load_data(CSV_PATH)
    df = clean_data(df)
    df = detect_anomalies(df)

    # Ensure numeric columns
    df['clicks'] = pd.to_numeric(df.get('clicks', 0), errors='coerce').fillna(0)
    df['impressions'] = pd.to_numeric(df.get('impressions', 0), errors='coerce').fillna(0)
    df['revenue'] = pd.to_numeric(df.get('revenue', 0), errors='coerce').fillna(0)

# Compute CTR safely (0 if no impressions)
    df['ctr'] = df.apply(lambda r: r['clicks'] / r['impressions'] if r['impressions'] > 0 else 0.0, axis=1)
 
    clicks, impressions, revenue, ctr = compute_kpis(df)

# Create main daily clicks chart (this uses df that may have anomaly column)
    chart_path = create_chart(df)

# Create CTR & Revenue charts
    ctr_chart = create_ctr_chart(df)
    revenue_chart = create_revenue_chart(df)

    basic_summary = create_summary(clicks, impressions, revenue, ctr, df)
    summary = ai_refine_summary(basic_summary)

# Pass all charts into PPTX creator
    create_pptx(summary, chart_path, ctr_chart=ctr_chart, revenue_chart=revenue_chart)


    print("\nInsightFlow is done! 🎉")
    print("Open your report here:", PPTX_PATH)

if __name__ == "__main__":
    main()
